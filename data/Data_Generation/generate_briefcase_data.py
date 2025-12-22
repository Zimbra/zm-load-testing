#!/usr/bin/env python3
import os
import random
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image, ImageDraw, ImageFont

# ---- Size & text config ----
MIN_SIZE_KB = 2000
MAX_SIZE_KB = 4000
TARGET_MB = 100
TARGET_BYTES = TARGET_MB * 1024 * 1024

SENTENCES = [
    "Artificial intelligence is transforming the world.",
    "Python is widely used for automation and data science.",
    "Machine learning models improve through data and iteration.",
    "Cybersecurity is essential for all modern organizations.",
    "Cloud computing allows access to resources over the internet.",
    "Innovation drives progress in technology and society.",
    "Natural language processing enables computers to understand text.",
    "Automation simplifies repetitive tasks and increases efficiency.",
    "Blockchain provides a secure way to record transactions.",
]

# ---- Utility helpers ----
def random_paragraph():
    return " ".join(random.choices(SENTENCES, k=random.randint(5, 10)))

def get_target_size_bytes():
    return random.randint(MIN_SIZE_KB, MAX_SIZE_KB) * 1024

# ---- File creators ----
def create_txt(file_path):
    target_size = get_target_size_bytes()
    content = ""
    while len(content.encode("utf-8")) < target_size:
        content += random_paragraph() + "\n"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)

def create_jpg(file_path):
    width, height = 1024, 768
    img = Image.new('RGB', (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    text = (random_paragraph() + " ") * 30
    draw.text((10, 10), text, fill=(0, 0, 0), font=font)
    img.save(file_path, quality=random.randint(60, 95))

def create_docx(file_path):
    doc = Document()
    doc.add_heading("Random Report", 0)
    for _ in range(random.randint(5, 15)):
        doc.add_paragraph(random_paragraph())
    doc.save(file_path)

def create_pptx(file_path):
    prs = Presentation()
    for _ in range(random.randint(2, 6)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Presentation Slide"
        slide.placeholders[1].text = random_paragraph()
    prs.save(file_path)

def create_xlsx(file_path):
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    for _ in range(50):
        ws.append([random.choice(SENTENCES) for _ in range(5)])
    wb.save(file_path)

def create_pdf(file_path):
    c = canvas.Canvas(file_path, pagesize=letter)
    text_obj = c.beginText(40, 750)
    text_obj.setFont("Helvetica", 12)
    text_obj.textLine("Generated PDF Report")
    for _ in range(50):
        text_obj.textLine(random_paragraph())
    c.drawText(text_obj)
    c.showPage()
    c.save()

# ---- Core generator ----
def generate_files_for_users(user_file="userlist.txt", output_dir="generated_files"):
    if not os.path.exists(user_file):
        raise FileNotFoundError(f"{user_file} not found!")

    with open(user_file, "r") as f:
        users = [line.strip() for line in f if line.strip()]

    os.makedirs(output_dir, exist_ok=True)

    creators = {
        "txt": create_txt,
        "jpg": create_jpg,
        "docx": create_docx,
        "pptx": create_pptx,
        "xlsx": create_xlsx,
        "pdf": create_pdf,
    }

    for user in users:
        user_dir = os.path.join(output_dir, user)
        os.makedirs(user_dir, exist_ok=True)

        total_bytes = 0
        file_count = 0
        print(f"\n📂 Generating files for {user}...")

        while total_bytes < TARGET_BYTES:
            ext = random.choice(list(creators.keys()))
            file_count += 1
            file_path = os.path.join(user_dir, f"{user}_{file_count}.{ext}")
            try:
                creators[ext](file_path)
            except Exception as e:
                print(f"[!] Error creating {ext}: {e}")
                continue

            current_size = os.path.getsize(file_path)
            total_bytes += current_size
            print(f"  + {os.path.basename(file_path)} ({current_size//1024} KB, total={total_bytes//1024} KB)", end="\r")

        total_mb = total_bytes / (1024 * 1024)
        print(f"\n✅ Done: {user} folder reached {total_mb:.2f} MB")

    print(f"\n🎉 Completed for {len(users)} users.")

# ---- Entry point ----
if __name__ == "__main__":
    generate_files_for_users()

