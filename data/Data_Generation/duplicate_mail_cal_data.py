#!/usr/bin/env python3
"""
duplicate_mail_cal_data.py

Each sender in from_users.txt sends messages to ALL recipients in to_users.txt.
dup-mail-mb and dup-cal-mb are targets per sender (sender mailbox only).

Features:
 - Deterministic fixed body sizes (configurable with --avg-msg-size)
 - Attachment percentage and weighted distribution (--attachment-percent, --attachment-type-dist)
 - Deterministic attachment pool (fixed per-type sizes)
 - Calendar messages contain a text/calendar MIME alternative (no separate .ics file)
 - Optional SMTP sending (provide --smtp-host to actually send)
 - Detailed per-sender report: Mail vs Calendar, With/Without attachments, attachment type breakdown
"""

from __future__ import annotations
import os
import sys
import argparse
import random
import time
import smtplib
import base64
from io import BytesIO
from email.message import EmailMessage
from email.utils import make_msgid, formatdate
from datetime import datetime, timedelta

# Third-party libs used in attachment generation:
# Pillow, python-docx, xlsxwriter, python-pptx, reportlab (if PDF needed)
try:
    from PIL import Image
    from docx import Document
    import xlsxwriter
    from pptx import Presentation
except Exception:
    # If these are missing you can still run but attachments for those types won't be generated.
    Image = None
    Document = None
    xlsxwriter = None
    Presentation = None

# -------------------------
# Helpers
# -------------------------
def parse_kv_dist(s: str):
    parts = [p.strip() for p in s.split(",") if p.strip()]
    if not parts:
        raise argparse.ArgumentTypeError("Empty distribution")
    out = []
    total = 0.0
    for part in parts:
        if "=" not in part:
            raise argparse.ArgumentTypeError(f"Bad part: {part}")
        k, v = part.split("=", 1)
        try:
            w = float(v.strip())
        except Exception:
            raise argparse.ArgumentTypeError(f"Bad weight for {k}: {v}")
        out.append((k.strip(), w))
        total += w
    if total <= 0:
        raise argparse.ArgumentTypeError("Distribution weights must sum to > 0")
    return out

def build_cumulative(kv_list):
    cum = []
    s = 0.0
    for k, w in kv_list:
        s += float(w)
        cum.append((s, k))
    return cum, s

def weighted_choice(cum_list, total_weight):
    r = random.uniform(0, total_weight)
    for threshold, key in cum_list:
        if r <= threshold:
            return key
    return cum_list[-1][1]

# -------------------------
# Deterministic body
# -------------------------
def create_fixed_body(size_bytes: int) -> bytes:
    """Create deterministic body of exactly size_bytes (approx)."""
    # Use ASCII 'A' to avoid multi-byte issues.
    return ("A" * size_bytes).encode("ascii", "ignore")

# -------------------------
# Attachment generators (simple deterministic, fixed-size)
# -------------------------
def create_jpeg_bytes(size_bytes=10 * 1024):
    if Image is None:
        return b"A" * size_bytes
    img = Image.new("RGB", (100, 100), (200, 100, 100))
    bio = BytesIO()
    img.save(bio, format="JPEG", quality=75)
    data = bio.getvalue()
    if len(data) < size_bytes:
        data += b"\0" * (size_bytes - len(data))
    else:
        data = data[:size_bytes]
    return data

def create_png_bytes(size_bytes=10 * 1024):
    if Image is None:
        return b"A" * size_bytes
    img = Image.new("RGB", (100, 100), (100, 200, 100))
    bio = BytesIO()
    img.save(bio, format="PNG")
    data = bio.getvalue()
    if len(data) < size_bytes:
        data += b"\0" * (size_bytes - len(data))
    else:
        data = data[:size_bytes]
    return data

def create_pdf_bytes(size_bytes=10 * 1024):
    # Lightweight deterministic PDF-like content (not real PDF library required)
    content = b"%PDF-1.4\n%Generated\n" + (b"A" * (size_bytes - 20))
    return content[:size_bytes]

def create_docx_bytes(size_bytes=10 * 1024):
    if Document is None:
        return b"A" * size_bytes
    doc = Document()
    doc.add_paragraph("Test DOCX")
    bio = BytesIO()
    doc.save(bio)
    data = bio.getvalue()
    if len(data) < size_bytes:
        data += b"\0" * (size_bytes - len(data))
    else:
        data = data[:size_bytes]
    return data

def create_xlsx_bytes(size_bytes=10 * 1024):
    if xlsxwriter is None:
        return b"A" * size_bytes
    bio = BytesIO()
    wb = xlsxwriter.Workbook(bio, {'in_memory': True})
    ws = wb.add_worksheet()
    ws.write(0, 0, "Test")
    wb.close()
    data = bio.getvalue()
    if len(data) < size_bytes:
        data += b"\0" * (size_bytes - len(data))
    else:
        data = data[:size_bytes]
    return data

def create_pptx_bytes(size_bytes=10 * 1024):
    if Presentation is None:
        return b"A" * size_bytes
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(0, 0, 100, 50)
    tb.text = "Test PPT"
    bio = BytesIO()
    prs.save(bio)
    data = bio.getvalue()
    if len(data) < size_bytes:
        data += b"\0" * (size_bytes - len(data))
    else:
        data = data[:size_bytes]
    return data

def create_txt_bytes(size_bytes=10 * 1024):
    return (b"A" * size_bytes)[:size_bytes]

# -------------------------
# Build attachment pool (fixed-size per-type)
# -------------------------
def build_attachment_pool(ext_list, per_type=20, fixed_kb=10):
    pool = {}
    size_bytes = fixed_kb * 1024
    for ext in ext_list:
        pool[ext] = []
        for i in range(per_type):
            if ext in ("jpg", "jpeg"):
                pool[ext].append(create_jpeg_bytes(size_bytes))
            elif ext == "png":
                pool[ext].append(create_png_bytes(size_bytes))
            elif ext == "pdf":
                pool[ext].append(create_pdf_bytes(size_bytes))
            elif ext == "docx":
                pool[ext].append(create_docx_bytes(size_bytes))
            elif ext in ("ppt", "pptx"):
                pool[ext].append(create_pptx_bytes(size_bytes))
            elif ext == "xlsx":
                pool[ext].append(create_xlsx_bytes(size_bytes))
            elif ext == "txt":
                pool[ext].append(create_txt_bytes(size_bytes))
            else:
                pool[ext].append(b"A" * size_bytes)
    return pool

# -------------------------
# Build a message (mail or calendar)
# -------------------------
def build_message(sender: str, recipients: list[str], is_calendar: bool,
                  attach_blob: bytes | None, attach_ext: str | None,
                  attachment_mimes: dict, body_bytes: bytes) -> EmailMessage:
    msg = EmailMessage()
    msg["From"] = sender
    msg["To"] = ", ".join(recipients)
    msg["Date"] = formatdate(localtime=True)
    msg["Message-ID"] = make_msgid()
    # Body
    msg.set_content(body_bytes.decode("utf-8", "ignore"))
    # Calendar alternative (text/calendar) for calendar messages
    if is_calendar:
        dtstart = datetime.utcnow() + timedelta(hours=1)
        dtend = dtstart + timedelta(hours=1)
        uid = make_msgid().strip("<>")
        ical = (
            "BEGIN:VCALENDAR\r\n"
            "VERSION:2.0\r\n"
            "PRODID:-//Duplicate Sender//EN\r\n"
            "METHOD:REQUEST\r\n"
            "BEGIN:VEVENT\r\n"
            f"UID:{uid}\r\n"
            f"DTSTAMP:{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}\r\n"
            f"DTSTART:{dtstart.strftime('%Y%m%dT%H%M%SZ')}\r\n"
            f"DTEND:{dtend.strftime('%Y%m%dT%H%M%SZ')}\r\n"
            f"SUMMARY:Generated Event\r\n"
            f"DESCRIPTION:Generated Event\r\n"
            "END:VEVENT\r\n"
            "END:VCALENDAR\r\n"
        )
        # Add calendar alternative
        msg.add_alternative(ical, subtype="calendar;method=REQUEST")
    # Attachment if present
    if attach_blob is not None and attach_ext is not None:
        maintype, subtype = attachment_mimes.get(attach_ext, ("application", "octet-stream"))
        filename = f"att_{attach_ext}"
        msg.add_attachment(attach_blob, maintype=maintype, subtype=subtype, filename=filename)
    return msg

# -------------------------
# SMTP send helper
# -------------------------
def smtp_send_raw(host: str, port: int, sender: str, recipients: list[str], raw_bytes: bytes, timeout=30):
    try:
        s = smtplib.SMTP(host, port, timeout=timeout)
        s.sendmail(sender, recipients, raw_bytes)
        s.quit()
    except Exception as e:
        print("[WARN] SMTP send failed:", e)

# -------------------------
# Main deterministic sender (Mode A1: sender -> ALL recipients; targets per sender)
# -------------------------
def run_deterministic_sender(args):
    # Load lists
    from_path = args.from_users
    to_path = args.to_users
    if not os.path.exists(from_path) or not os.path.exists(to_path):
        print("ERROR: from-users or to-users file not found.")
        sys.exit(1)

    from_users = [l.strip() for l in open(from_path, "r").read().splitlines() if l.strip()]
    to_users = [l.strip() for l in open(to_path, "r").read().splitlines() if l.strip()]

    if not from_users or not to_users:
        print("ERROR: from-users or to-users is empty.")
        sys.exit(1)

    # parse attachment distribution
    att_kv = parse_kv_dist(args.attachment_type_dist)
    att_cum, att_total = build_cumulative(att_kv)
    att_keys = [k for k,_ in att_kv]

    # Build deterministic attachment pool
    att_pool = build_attachment_pool(att_keys, per_type=20, fixed_kb=args.attachment_kb)

    # attachment mime mapping
    attachment_mimes = {
        "pdf": ("application", "pdf"),
        "docx": ("application", "vnd.openxmlformats-officedocument.wordprocessingml.document"),
        "txt": ("text", "plain"),
        "png": ("image", "png"),
        "jpg": ("image", "jpeg"),
        "ppt": ("application", "vnd.ms-powerpoint"),
        "xlsx": ("application", "vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    }

    # Targets (per sender)
    mail_target = args.dup_mail_mb * 1024 * 1024
    cal_target = args.dup_cal_mb * 1024 * 1024

    # Track bytes (per sender). For reporting we also track recipients bytes separately.
    bytes_mail = {u: 0 for u in from_users}
    bytes_cal = {u: 0 for u in from_users}
    bytes_received = {u: 0 for u in set(from_users + to_users)}

    # Report structure (per sender)
    report = {}
    for u in from_users:
        report[u] = {
            "mail": {"with_att": 0, "without_att": 0, "types": {k: 0 for k in att_keys}},
            "calendar": {"with_att": 0, "without_att": 0, "types": {k: 0 for k in att_keys}}
        }

    # Round-robin index per attachment type to pick a file from the pool deterministically
    att_index = {k: 0 for k in att_keys}
    global_index = 0

    total_msgs = 0

    print(f"Starting deterministic duplicate send (Mode A1). Senders: {len(from_users)}, Recipients: {len(to_users)}")
    # Main loop: iterate senders until each sender meets both mail and calendar targets
    while True:
        all_done = True
        for sender in from_users:
            # Determine whether this sender still needs mail or calendar
            need_mail = bytes_mail[sender] < mail_target
            need_cal = bytes_cal[sender] < cal_target

            # If sender already satisfied both, skip
            if not (need_mail or need_cal):
                continue

            all_done = False

            # Priority: generate mail messages first until mail_target reached, then calendar
            if need_mail:
                is_calendar = False
            elif need_cal:
                is_calendar = True
            else:
                # shouldn't reach here
                continue

            # Decide whether this message will have an attachment
            has_attachment = (random.random() < (args.attachment_percent / 100.0))

            attach_blob = None
            attach_ext = None
            if has_attachment:
                # choose type according to distribution
                ext = weighted_choice(att_cum, att_total)
                # round-robin pick from pool for that ext
                idx = att_index[ext] % len(att_pool[ext])
                attach_blob = att_pool[ext][idx]
                attach_ext = ext
                att_index[ext] += 1

            # Fixed body size based on avg_msg_size
            body_bytes = create_fixed_body(args.avg_msg_size * 1024)

            # Build message
            msg = build_message(sender, to_users, is_calendar, attach_blob, attach_ext, attachment_mimes, body_bytes)
            raw = msg.as_bytes()
            size = len(raw)

            # Optionally send via SMTP
            if args.smtp_host:
                smtp_send_raw(args.smtp_host, args.smtp_port, sender, to_users, raw)

            # Track bytes (per sender goal)
            if is_calendar:
                bytes_cal[sender] += size
            else:
                bytes_mail[sender] += size

            # Track recipients bytes for reporting (but NOT used to stop loop)
            bytes_received[sender] += size
            for r in to_users:
                bytes_received[r] += size

            # Update report per sender
            key = "calendar" if is_calendar else "mail"
            if attach_blob:
                report[sender][key]["with_att"] += size
                report[sender][key]["types"][attach_ext] += size
            else:
                report[sender][key]["without_att"] += size

            global_index += 1
            total_msgs += 1

        if all_done:
            break

    # Print final report
    print(f"\nCompleted. Total messages generated: {total_msgs}\n")
    print("=== Duplicate Data Report (per sender) ===")
    for sender in from_users:
        print(f"\nSender: {sender}")
        mail_with = report[sender]["mail"]["with_att"] / (1024 * 1024)
        mail_wo = report[sender]["mail"]["without_att"] / (1024 * 1024)
        print(f"  Mail:   With attachments: {mail_with:.2f} MB   Without attachments: {mail_wo:.2f} MB   Total: {(mail_with+mail_wo):.2f} MB")
        if any(report[sender]["mail"]["types"].values()):
            print("    Mail attachment breakdown:")
            for ext, b in report[sender]["mail"]["types"].items():
                if b:
                    print(f"      {ext.upper():6}: {b/(1024*1024):.2f} MB")

        cal_with = report[sender]["calendar"]["with_att"] / (1024 * 1024)
        cal_wo = report[sender]["calendar"]["without_att"] / (1024 * 1024)
        print(f"  Calendar: With attachments: {cal_with:.2f} MB   Without attachments: {cal_wo:.2f} MB   Total: {(cal_with+cal_wo):.2f} MB")
        if any(report[sender]["calendar"]["types"].values()):
            print("    Calendar attachment breakdown:")
            for ext, b in report[sender]["calendar"]["types"].items():
                if b:
                    print(f"      {ext.upper():6}: {b/(1024*1024):.2f} MB")

    print("\nDone.")

# -------------------------
# CLI
# -------------------------
def main():
    p = argparse.ArgumentParser(description="Deterministic duplicate sender (Mode A1: each sender -> ALL recipients)")
    p.add_argument("--from-users", required=True, help="Path to from_users.txt (one sender per line)")
    p.add_argument("--to-users", required=True, help="Path to to_users.txt (one recipient per line)")
    p.add_argument("--dup-mail-mb", type=int, required=True, help="Per-sender target duplicate mail MB")
    p.add_argument("--dup-cal-mb", type=int, default=0, help="Per-sender target duplicate calendar MB")
    p.add_argument("--attachment-percent", type=float, default=50.0, help="Percent of messages that include attachment")
    p.add_argument("--attachment-type-dist", type=str, default="pdf=30,docx=20,txt=10,png=10,jpg=10,ppt=10,xlsx=10",
                   help="Weighted distribution for attachment types (csv: ext=weight,...)")
    p.add_argument("--attachment-kb", type=int, default=10, help="Fixed size (KB) for generated attachments (default 10KB)")
    p.add_argument("--avg-msg-size", type=int, default=50, help="Fixed message body size in KB (default 50KB)")
    p.add_argument("--smtp-host", type=str, default=None, help="SMTP host to send messages (optional)")
    p.add_argument("--smtp-port", type=int, default=25, help="SMTP port (default 25)")
    args = p.parse_args()
    run_deterministic_sender(args)

if __name__ == "__main__":
    main()
